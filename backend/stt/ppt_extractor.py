import argparse
import json
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from .image_description import describe_image


def extract_pptx(path: str, describe_images: bool = True) -> dict:
    """describe_images=True calls a vision LLM (image_description.describe_image, costs
    money) for every picture shape found — set False to skip images entirely (old
    behavior, text-only, free)."""
    prs = Presentation(path)
    slides = []
    for i, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in paragraph.runs).strip()
                    if text:
                        texts.append(text)
            elif describe_images and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                description = describe_image(shape.image.blob, shape.image.content_type)
                texts.append(f"[이미지 설명: {description}]")

        notes = ""
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()

        slides.append({
            "slide_number": i,
            "text": "\n".join(texts),
            "notes": notes,
        })

    return {"source": Path(path).name, "slides": slides}


def main():
    parser = argparse.ArgumentParser(description="Extract slide text from a .pptx file")
    parser.add_argument("pptx_path")
    parser.add_argument(
        "--no-image-description",
        action="store_true",
        help="Skip vision-LLM image descriptions (text-only, free, old behavior)",
    )
    parser.add_argument("-o", "--output", help="Write JSON to this path instead of stdout")
    args = parser.parse_args()

    result = extract_pptx(args.pptx_path, describe_images=not args.no_image_description)
    output = json.dumps(result, ensure_ascii=False, indent=2)

    if args.output:
        Path(args.output).write_text(output, encoding="utf-8")
    else:
        print(output)


if __name__ == "__main__":
    main()
